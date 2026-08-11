"""Generate synthetic test fixtures.

The fixtures are deliberately fictional so the test suite never depends on a
real, confidential pitch deck. Run with::

    python tests/fixtures/make_fixtures.py
"""

from __future__ import annotations

from pathlib import Path

FIXTURE_DIR = Path(__file__).parent

#: (title, [body lines]) - a small but realistic fictional seed-stage SaaS deck.
SLIDES: list[tuple[str, list[str]]] = [
    (
        "Meridian Freight OS",
        [
            "Exception management for mid-market freight brokers",
            "Seed round - March 2026",
            "hello@meridianfreight.example",
        ],
    ),
    (
        "The problem",
        [
            "• Mid-market brokers run 60-80% of shipment exceptions through email and phone",
            "• A single mishandled exception costs an average of $410 in penalties and re-work",
            "• Brokers we surveyed lose 11 hours per dispatcher per week to exception chasing",
            "• Source: internal survey of 42 brokerages, Q4 2025",
        ],
    ),
    (
        "Why now",
        [
            "• FMCSA broker transparency rules took effect January 2026",
            "• Shippers now demand exception audit trails as a contract term",
            "• Legacy TMS vendors price integrations at $60K+ per year, out of reach mid-market",
        ],
    ),
    (
        "Product",
        [
            "• Meridian ingests EDI 214 and email status updates and classifies exceptions automatically",
            "• Dispatcher console proposes a resolution and drafts the customer update",
            "• Deploys in 3 weeks alongside the incumbent TMS - no rip and replace",
            "• Live integrations: McLeod, Turvo, Aljex",
        ],
    ),
    (
        "Traction",
        [
            "• $2.4M ARR as of Q2 2026",
            "• 118% year-over-year ARR growth",
            "• 37 paying brokerages",
            "• Gross margin 71%",
            "• Net revenue retention 114%",
            "• Average contract value $64,800",
            "• Logo churn 6% over the trailing twelve months",
        ],
    ),
    (
        "Customer evidence",
        [
            "• Ridgeline Logistics cut exception handling time 42% in the first 90 days",
            "• Two of the top five customers expanded seats within 6 months",
            "• Sales cycle averages 74 days from first call to signature",
            "• 3 of 37 customers are on annual contracts; the rest are monthly",
        ],
    ),
    (
        "Market",
        [
            "• 17,000 licensed freight brokerages in the United States",
            "• We target the 4,100 brokerages doing $10M-$250M in annual freight volume",
            "• At our current $64,800 ACV that is a $266M serviceable market",
            "• Source: FMCSA licensing data and company analysis",
        ],
    ),
    (
        "Business model",
        [
            "• Annual SaaS subscription priced per dispatcher seat",
            "• Land at 8-12 seats, expand to 30+ as brokerages consolidate desks",
            "• CAC $19,400 blended; payback in 14 months at current gross margin",
        ],
    ),
    (
        "Competition",
        [
            "• Incumbent TMS vendors (McLeod, Turvo) bundle basic status tracking",
            "• Point solutions focus on tracking visibility, not exception resolution",
            "• Our exception classification model is trained on 2.1M labelled shipment events",
        ],
    ),
    (
        "Team",
        [
            "• Dana Whitfield, CEO - 9 years at C.H. Robinson, ran a 60-person brokerage desk",
            "• Amir Kohli, CTO - built dispatch systems at Convoy, 11 years in logistics software",
            "• 14 full-time employees; 4 in engineering, 3 in customer success",
            "• No full-time VP of Sales; the CEO closes all enterprise deals today",
        ],
    ),
    (
        "The ask",
        [
            "• Raising $6M Seed",
            "• $12M pre-money valuation cap",
            "• Use of proceeds: 55% engineering, 30% go-to-market, 15% general and administrative",
            "• 24 months of runway to $10M ARR",
        ],
    ),
    (
        "Milestones",
        [
            "• Q4 2026: 75 customers and $5M ARR",
            "• Q2 2027: self-serve onboarding live, CAC below $12,000",
            "• Q4 2027: $10M ARR and first two enterprise shipper contracts",
        ],
    ),
]


def build_pdf(path: Path) -> Path:
    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.pdfgen.canvas import Canvas

    width, height = landscape(letter)
    canvas = Canvas(str(path), pagesize=(width, height))

    for title, lines in SLIDES:
        canvas.setFont("Helvetica-Bold", 26)
        canvas.drawString(60, height - 90, title)
        canvas.setFont("Helvetica", 15)
        y = height - 150
        for line in lines:
            canvas.drawString(60, y, line)
            y -= 28
        canvas.setFont("Helvetica", 8)
        canvas.drawString(60, 36, "Meridian Freight OS - Confidential")
        canvas.showPage()

    canvas.save()
    return path


def build_pptx(path: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    presentation = Presentation()
    blank = presentation.slide_layouts[6]

    for title, lines in SLIDES:
        slide = presentation.slides.add_slide(blank)

        title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(8.8), Inches(0.9))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].runs[0].font.size = Pt(30)
        title_frame.paragraphs[0].runs[0].font.bold = True

        body_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.6), Inches(8.8), Inches(4.5))
        body_frame = body_box.text_frame
        body_frame.word_wrap = True
        for index, line in enumerate(lines):
            paragraph = body_frame.paragraphs[0] if index == 0 else body_frame.add_paragraph()
            paragraph.text = line.lstrip("• ").strip()
            paragraph.runs[0].font.size = Pt(16)

        slide.notes_slide.notes_text_frame.text = f"Speaker notes for '{title}'."

    presentation.save(str(path))
    return path


def main() -> None:
    FIXTURE_DIR.mkdir(parents=True, exist_ok=True)
    pdf = build_pdf(FIXTURE_DIR / "sample_pitch.pdf")
    pptx = build_pptx(FIXTURE_DIR / "sample_pitch.pptx")
    print(f"Wrote {pdf}")
    print(f"Wrote {pptx}")


if __name__ == "__main__":
    main()
