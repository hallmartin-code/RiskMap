"""PPTX ingestion via python-pptx.

Shapes are read in top-to-bottom, left-to-right order so the extracted text
follows the visual reading order of the slide rather than XML order.
"""

from __future__ import annotations

from pathlib import Path

from ..errors import CorruptDocumentError
from ..extraction.metric_extractor import extract_metrics
from ..extraction.text_cleaner import clean_page_text, split_bullets, strip_repeated_lines
from ..logging_setup import get_logger
from ..models.deck import DeckDocument, DeckPage, DeckTable

log = get_logger("ingestion.pptx")

IMAGE_ONLY_CHAR_THRESHOLD = 40


def _shape_sort_key(shape) -> tuple[int, int]:
    top = getattr(shape, "top", None)
    left = getattr(shape, "left", None)
    return (top if top is not None else 10**9, left if left is not None else 10**9)


def _iter_text_shapes(shapes):
    """Yield shapes in reading order, descending into groups."""
    for shape in sorted(shapes, key=_shape_sort_key):
        if getattr(shape, "shape_type", None) is not None and shape.shape_type == 6:  # GROUP
            yield from _iter_text_shapes(shape.shapes)
        else:
            yield shape


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    lines: list[str] = []
    for paragraph in shape.text_frame.paragraphs:
        text = "".join(run.text for run in paragraph.runs).strip()
        if not text:
            continue
        # Preserve list structure so bullets survive cleanup.
        lines.append(f"• {text}" if paragraph.level > 0 else text)
    return "\n".join(lines)


def _table_rows(shape) -> list[list[str]]:
    rows: list[list[str]] = []
    for row in shape.table.rows:
        rows.append([cell.text.strip() for cell in row.cells])
    return [r for r in rows if any(r)]


def parse_pptx(path: Path, original_name: str | None = None, file_type: str = "pptx") -> DeckDocument:
    """Parse a .pptx deck into the normalised representation."""
    from pptx import Presentation

    try:
        presentation = Presentation(str(path))
    except Exception as exc:
        raise CorruptDocumentError(f"Could not open PowerPoint file '{path.name}': {exc}") from exc

    warnings: list[str] = []
    raw_pages: list[str] = []
    slide_titles: list[str | None] = []
    slide_tables: list[list[list[list[str]]]] = []
    slide_notes: list[str] = []

    for index, slide in enumerate(presentation.slides, start=1):
        title: str | None = None
        try:
            if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
                title = slide.shapes.title.text.strip() or None
        except Exception:  # pragma: no cover - unusual layouts
            title = None

        chunks: list[str] = []
        tables: list[list[list[str]]] = []
        for shape in _iter_text_shapes(slide.shapes):
            try:
                if getattr(shape, "has_table", False):
                    rows = _table_rows(shape)
                    if rows:
                        tables.append(rows)
                    continue
                text = _shape_text(shape)
                if text and text.strip() != (title or "").strip():
                    chunks.append(text)
            except Exception as exc:  # pragma: no cover
                warnings.append(f"Slide {index}: shape skipped ({exc}).")

        notes = ""
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:  # pragma: no cover
            notes = ""

        raw_pages.append("\n".join(chunks))
        slide_titles.append(title)
        slide_tables.append(tables)
        slide_notes.append(notes)

    deduped = strip_repeated_lines(raw_pages)

    pages: list[DeckPage] = []
    for index, raw in enumerate(deduped, start=1):
        text = clean_page_text(raw)
        tables = [DeckTable(page_number=index, rows=rows) for rows in slide_tables[index - 1]]
        searchable = "\n".join([text, "\n".join(t.to_text() for t in tables), slide_notes[index - 1]])
        title = slide_titles[index - 1]
        body_len = len(text) + len(slide_notes[index - 1]) + sum(len(t.to_text()) for t in tables)

        pages.append(
            DeckPage(
                page_number=index,
                title=title,
                raw_text=text,
                bullets=split_bullets(text),
                tables=tables,
                metrics=extract_metrics(searchable, index),
                speaker_notes=slide_notes[index - 1],
                is_image_only=body_len < IMAGE_ONLY_CHAR_THRESHOLD and not title,
            )
        )

    log.info("Parsed PPTX '%s': %d slides", original_name or path.name, len(pages))
    return DeckDocument(
        filename=original_name or path.name,
        file_type=file_type,  # type: ignore[arg-type]
        pages=pages,
        warnings=warnings,
    )
